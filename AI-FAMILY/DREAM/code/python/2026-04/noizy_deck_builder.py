#!/usr/bin/env python3
"""
NOIZY.AI EMPIRE STRATEGY DECK
Master Identity System Architect — python-pptx generator

Run:
    pip3 install python-pptx pillow
    python3 noizy_deck_builder.py

Output: noizy_empire_deck.pptx (same directory as this script)

Author: RSP / MC96ECO — God Node M2 Ultra
"""

import os
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
    from pptx.enum.dml import MSO_THEME_COLOR
    import pptx.oxml.ns as oxml_ns
    from lxml import etree
except ImportError:
    print("Installing dependencies...")
    os.system(f"{sys.executable} -m pip install python-pptx pillow --quiet")
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
    from lxml import etree

# ── PALETTE ──────────────────────────────────────────────────────────────────
BG_DARK      = RGBColor(0x08, 0x08, 0x10)   # near-black, deep navy
BG_CARD      = RGBColor(0x0F, 0x0F, 0x1E)   # card surface
ACCENT_CYAN  = RGBColor(0x00, 0xE5, 0xFF)   # electric cyan — main brand
ACCENT_GOLD  = RGBColor(0xFF, 0xC2, 0x00)   # sovereing gold
ACCENT_PINK  = RGBColor(0xFF, 0x2D, 0x78)   # hot pink / danger
ACCENT_GREEN = RGBColor(0x39, 0xFF, 0x8E)   # neon green / go
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
GREY_MID     = RGBColor(0x8A, 0x8A, 0xAA)
GREY_DARK    = RGBColor(0x1E, 0x1E, 0x32)

# ── SLIDE DIMENSIONS (16:9 widescreen) ───────────────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

OUT_FILE = Path(__file__).parent / "noizy_empire_deck.pptx"


# ── HELPERS ───────────────────────────────────────────────────────────────────

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(blank_layout)


def fill_bg(slide, color: RGBColor = BG_DARK):
    """Fill the entire slide background with a solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height,
             fill_color=None, line_color=None, line_width_pt=1.0):
    """Add a filled/outlined rectangle shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    if fill_color:
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()   # transparent

    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Arial"
    return txBox


def add_label_value(slide, label, value, left, top, width,
                    label_color=ACCENT_CYAN, value_color=WHITE,
                    label_size=11, value_size=14):
    """Add a label + value stacked pair."""
    add_text(slide, label.upper(), left, top, width, Inches(0.3),
             font_size=label_size, bold=True, color=label_color)
    add_text(slide, value, left, top + Inches(0.28), width, Inches(0.4),
             font_size=value_size, color=value_color)


def add_divider(slide, left, top, width, color=ACCENT_CYAN, thickness=1.5):
    """Thin horizontal rule."""
    shape = slide.shapes.add_shape(1, left, top, width, Pt(thickness))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_tag(slide, text, left, top, fill=ACCENT_CYAN, text_color=BG_DARK, size=9):
    """Small pill tag."""
    w = Inches(len(text) * 0.085 + 0.25)
    h = Inches(0.22)
    rect = add_rect(slide, left, top, w, h, fill_color=fill)
    add_text(slide, text, left, top, w, h,
             font_size=size, bold=True, color=text_color, align=PP_ALIGN.CENTER)
    return w


def slide_header(slide, title, subtitle=None,
                 icon="◈", accent=ACCENT_CYAN, slide_num=None):
    """Standard slide header with top accent bar, icon, title, subtitle."""

    # top accent bar
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), fill_color=accent)

    # slide number badge (top right)
    if slide_num:
        add_text(slide, f"{slide_num:02d}", SLIDE_W - Inches(0.8),
                 Inches(0.12), Inches(0.6), Inches(0.3),
                 font_size=9, color=GREY_MID, align=PP_ALIGN.RIGHT)

    # icon + title
    add_text(slide, icon, Inches(0.45), Inches(0.22), Inches(0.55), Inches(0.6),
             font_size=28, bold=True, color=accent)
    add_text(slide, title, Inches(0.95), Inches(0.22), Inches(11.0), Inches(0.55),
             font_size=30, bold=True, color=WHITE)

    # subtitle
    if subtitle:
        add_text(slide, subtitle, Inches(0.95), Inches(0.72),
                 Inches(10.0), Inches(0.4),
                 font_size=13, color=GREY_MID, italic=True)

    # bottom divider
    add_divider(slide, Inches(0.45), Inches(1.1), SLIDE_W - Inches(0.9), color=accent)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════

def slide_cover(prs):
    slide = blank_slide(prs)
    fill_bg(slide)

    # Full-width gradient bar (simulated with stacked thin rects)
    colors_h = [ACCENT_CYAN, RGBColor(0x00, 0xA0, 0xCC), RGBColor(0x00, 0x60, 0x80)]
    for i, c in enumerate(colors_h):
        add_rect(slide, 0, Inches(i * 0.55), SLIDE_W, Inches(0.55), fill_color=c)

    # Dark overlay on top 2/3
    overlay = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(5.5))
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = BG_DARK
    overlay.line.fill.background()

    # Brand tag line
    add_text(slide, "CLASSIFIED // INTERNAL STRATEGY BRIEF",
             Inches(0.5), Inches(0.6), Inches(10), Inches(0.35),
             font_size=10, bold=True, color=ACCENT_PINK, align=PP_ALIGN.LEFT)

    # Primary title
    add_text(slide, "noizy.ai", Inches(0.5), Inches(1.2), SLIDE_W - Inches(1),
             Inches(1.4), font_size=88, bold=True, color=ACCENT_CYAN,
             align=PP_ALIGN.LEFT)

    # Sub-title
    add_text(slide, "EMPIRE STRATEGY DECK",
             Inches(0.55), Inches(2.5), Inches(9), Inches(0.7),
             font_size=26, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    add_text(slide, "Master Identity System Architecture  ·  Agentic Mesh  ·  The Gathering  ·  Heaven Infrastructure",
             Inches(0.55), Inches(3.1), Inches(10), Inches(0.4),
             font_size=12, color=GREY_MID, align=PP_ALIGN.LEFT)

    # Metrics row
    metrics = [
        ("6", "Sub-Brands"),
        ("8+", "Integrations"),
        ("12TB+", "Asset Vault"),
        ("April 17", "Launch"),
    ]
    for i, (val, label) in enumerate(metrics):
        x = Inches(0.5 + i * 3.1)
        add_rect(slide, x, Inches(4.0), Inches(2.8), Inches(1.3),
                 fill_color=GREY_DARK, line_color=ACCENT_CYAN, line_width_pt=0.8)
        add_text(slide, val, x + Inches(0.15), Inches(4.12), Inches(2.5),
                 Inches(0.6), font_size=32, bold=True, color=ACCENT_CYAN,
                 align=PP_ALIGN.LEFT)
        add_text(slide, label, x + Inches(0.15), Inches(4.68), Inches(2.5),
                 Inches(0.3), font_size=11, color=GREY_MID, align=PP_ALIGN.LEFT)

    # Bottom bar
    add_rect(slide, 0, SLIDE_H - Inches(0.55), SLIDE_W, Inches(0.55),
             fill_color=GREY_DARK)
    add_text(slide, "RSP / MC96ECO  ·  God Node M2 Ultra  ·  HVS 75/25 Perpetual  ·  Gabriel Watching",
             Inches(0.5), SLIDE_H - Inches(0.48), SLIDE_W - Inches(1), Inches(0.4),
             font_size=9, color=GREY_MID, align=PP_ALIGN.LEFT)
    add_text(slide, "INTERNAL USE ONLY", SLIDE_W - Inches(2.5),
             SLIDE_H - Inches(0.48), Inches(2.2), Inches(0.4),
             font_size=9, bold=True, color=ACCENT_PINK, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — UNIFIED IDENTITY (SSO)
# ══════════════════════════════════════════════════════════════════════════════

def slide_sso(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    slide_header(slide, "The Unified Identity — noizy.ai SSO",
                 subtitle="Eliminating Dumb Tokens. One OAuth 2.0 flow to rule six brands.",
                 icon="◉", accent=ACCENT_CYAN, slide_num=2)

    # Central node
    cx, cy = Inches(6.6), Inches(3.8)
    node_w, node_h = Inches(2.4), Inches(1.0)
    add_rect(slide, cx - node_w/2, cy - node_h/2, node_w, node_h,
             fill_color=ACCENT_CYAN, line_color=WHITE, line_width_pt=1.5)
    add_text(slide, "noizy.ai", cx - node_w/2, cy - node_h/2, node_w, node_h,
             font_size=20, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)

    # Spoke nodes
    spokes = [
        (Inches(1.2),  Inches(2.2),  "Google Workspace",  "OAuth 2.0 IdP",    ACCENT_GOLD),
        (Inches(1.2),  Inches(4.8),  "Cloudflare",        "Workers + D1 + KV", ACCENT_CYAN),
        (Inches(5.6),  Inches(1.3),  "Stripe",            "Payments + HVS",    ACCENT_GREEN),
        (Inches(5.6),  Inches(6.0),  "Notion",            "Wisdom Archive",    ACCENT_GOLD),
        (Inches(10.5), Inches(2.2),  "Apple Dev Program", "iOS + iPadOS",      ACCENT_PINK),
        (Inches(10.5), Inches(4.8),  "Linear / GitHub",   "Ops + Source",      GREY_MID),
    ]

    for sx, sy, name, sub, color in spokes:
        sw, sh = Inches(2.0), Inches(0.85)
        add_rect(slide, sx, sy, sw, sh,
                 fill_color=GREY_DARK, line_color=color, line_width_pt=1.2)
        add_text(slide, name, sx + Inches(0.1), sy + Inches(0.06),
                 sw - Inches(0.15), Inches(0.4),
                 font_size=12, bold=True, color=color)
        add_text(slide, sub, sx + Inches(0.1), sy + Inches(0.45),
                 sw - Inches(0.15), Inches(0.35),
                 font_size=9, color=GREY_MID)

    # Key principle box
    add_rect(slide, Inches(0.45), Inches(6.1), Inches(12.4), Inches(0.85),
             fill_color=GREY_DARK, line_color=ACCENT_GOLD, line_width_pt=0.8)
    add_text(slide, "⚡  PRINCIPLE:  Single Google Workspace identity → cascaded permissions → all 6 brands → "
             "zero per-service signup friction → full Gabriel audit trail on every auth event.",
             Inches(0.65), Inches(6.2), Inches(12.0), Inches(0.65),
             font_size=11, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — AGENTIC MESH (THE TRINITY)
# ══════════════════════════════════════════════════════════════════════════════

def slide_mesh(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    slide_header(slide, "The Agentic Mesh — The Trinity",
                 subtitle="Gabriel routes. Lucy commands. Keith listens. Mickey P never drops a frame.",
                 icon="⬡", accent=ACCENT_PINK, slide_num=3)

    agents = [
        {
            "name":   "GABRIEL",
            "role":   "High-Level Router / Orchestrator",
            "detail": "Constitutional audit trail. Every event logged to D1 agent-memory. "
                      "Routes /api/dispatch → Mesh GOD:9696. HVS guardian.",
            "stack":  "Cloudflare Worker · D1 · KV · Durable Objects",
            "color":  ACCENT_GOLD,
            "tag":    "ORCHESTRATOR",
            "x":      Inches(0.45),
        },
        {
            "name":   "LUCY",
            "role":   "iPad UI & Spatial Controller",
            "detail": "Swift / React front-end. iPad command-and-control dashboard. "
                      "Spatial sessions over Cloudflare Access. Real-time Lucy ↔ Mesh sync.",
            "stack":  "Swift · React · CF Access · WebRTC (NOIZYSTREAM v2)",
            "color":  ACCENT_CYAN,
            "tag":    "UI / SPATIAL",
            "x":      Inches(4.72),
        },
        {
            "name":   "KEITH",
            "role":   "Audio Intelligence",
            "detail": "Voice-to-TypeScript pipeline. Whisper transcription at edge. "
                      "Emotional tag extraction. Feeds voice_profiles into Gabriel ledger.",
            "stack":  "@cf/openai/whisper · XTTS v2 · NOUIZYSTREAM v2",
            "color":  ACCENT_GREEN,
            "tag":    "AUDIO AI",
            "x":      Inches(8.99),
        },
    ]

    for ag in agents:
        x = ag["x"]
        w = Inches(4.0)
        # card
        add_rect(slide, x, Inches(1.3), w, Inches(5.5),
                 fill_color=GREY_DARK, line_color=ag["color"], line_width_pt=1.5)
        # top stripe
        add_rect(slide, x, Inches(1.3), w, Inches(0.08), fill_color=ag["color"])
        # tag
        add_tag(slide, ag["tag"], x + Inches(0.2), Inches(1.45),
                fill=ag["color"], text_color=BG_DARK, size=9)
        # name
        add_text(slide, ag["name"], x + Inches(0.2), Inches(1.72), w - Inches(0.4),
                 Inches(0.65), font_size=30, bold=True, color=ag["color"])
        # role
        add_text(slide, ag["role"], x + Inches(0.2), Inches(2.35), w - Inches(0.4),
                 Inches(0.4), font_size=12, italic=True, color=WHITE)
        # divider
        add_divider(slide, x + Inches(0.2), Inches(2.75), w - Inches(0.4),
                    color=ag["color"], thickness=0.8)
        # detail
        add_text(slide, ag["detail"], x + Inches(0.2), Inches(2.88), w - Inches(0.4),
                 Inches(2.0), font_size=11, color=GREY_MID)
        # stack label
        add_text(slide, "STACK", x + Inches(0.2), Inches(5.15), w - Inches(0.4),
                 Inches(0.25), font_size=8, bold=True, color=ag["color"])
        add_text(slide, ag["stack"], x + Inches(0.2), Inches(5.38), w - Inches(0.4),
                 Inches(0.6), font_size=9, color=WHITE)

    # Mickey P footnote
    add_rect(slide, Inches(0.45), Inches(6.98), Inches(12.4), Inches(0.42),
             fill_color=GREY_DARK)
    add_text(slide, "🎛️  MICKEY P  —  MacBook Pro audio routing daemon. Low-latency comms bridge. "
             "Never drops a frame. Feeds Keith's real-time pipeline.",
             Inches(0.65), Inches(7.02), Inches(12.0), Inches(0.35),
             font_size=10, color=GREY_MID)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — THE GATHERING
# ══════════════════════════════════════════════════════════════════════════════

def slide_gathering(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    slide_header(slide, '"The Gathering" — Asset Vacuum',
                 subtitle="12TB+ of fragmented genius. One repo. Zero data left behind.",
                 icon="◈", accent=ACCENT_GOLD, slide_num=4)

    # Left: process flow
    steps = [
        (ACCENT_GOLD,  "SCAN",    "Auto-scan M2 Ultra RAID + MAG drives for all code, audio, video, AI artifacts"),
        (ACCENT_CYAN,  "CLASSIFY","Tag by brand, date, type. Emotional tags applied by Keith. Sovereign C2PA stamps."),
        (ACCENT_GREEN, "COMMIT",  "Push to github.com/rsp-noizy/the-gathering — single source of truth"),
        (ACCENT_PINK,  "INDEX",   "Gabriel logs every asset. Notion wisdom archive updated. Linear tickets auto-created."),
    ]

    for i, (color, step, desc) in enumerate(steps):
        y = Inches(1.4 + i * 1.35)
        # number bubble (simulated)
        add_rect(slide, Inches(0.45), y, Inches(0.55), Inches(0.55),
                 fill_color=color)
        add_text(slide, str(i + 1), Inches(0.45), y, Inches(0.55), Inches(0.55),
                 font_size=18, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
        add_text(slide, step, Inches(1.1), y, Inches(2.5), Inches(0.3),
                 font_size=13, bold=True, color=color)
        add_text(slide, desc, Inches(1.1), y + Inches(0.3), Inches(5.8), Inches(0.75),
                 font_size=10, color=GREY_MID)
        # connector line (except last)
        if i < len(steps) - 1:
            add_rect(slide, Inches(0.69), y + Inches(0.55),
                     Pt(2), Inches(0.80), fill_color=GREY_MID)

    # Right: stats cards
    stats = [
        ("12 TB+",  "Total asset volume across all drives",    ACCENT_GOLD),
        ("6",       "Brand namespaces indexed",                 ACCENT_CYAN),
        ("1 repo",  "github.com/rsp-noizy/the-gathering",      ACCENT_GREEN),
        ("∞",       "Historical R&D preserved, never deleted",  ACCENT_PINK),
    ]
    for i, (val, label, color) in enumerate(stats):
        x = Inches(7.5 + (i % 2) * 2.9)
        y = Inches(1.4 + (i // 2) * 2.0)
        add_rect(slide, x, y, Inches(2.6), Inches(1.65),
                 fill_color=GREY_DARK, line_color=color, line_width_pt=1.2)
        add_text(slide, val, x + Inches(0.15), y + Inches(0.18), Inches(2.3),
                 Inches(0.75), font_size=36, bold=True, color=color)
        add_text(slide, label, x + Inches(0.15), y + Inches(0.9), Inches(2.3),
                 Inches(0.65), font_size=9, color=GREY_MID)

    # Bottom CTA
    add_rect(slide, Inches(0.45), Inches(6.65), Inches(12.4), Inches(0.72),
             fill_color=GREY_DARK, line_color=ACCENT_GOLD, line_width_pt=0.8)
    add_text(slide, "🎯  DESTINATION:  github.com/rsp-noizy/the-gathering  "
             "—  Sovereign. Immutable. Gabriel-audited. Every commit C2PA stamped.",
             Inches(0.65), Inches(6.75), Inches(12.0), Inches(0.55),
             font_size=12, bold=False, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — HEAVEN WORKER INFRASTRUCTURE
# ══════════════════════════════════════════════════════════════════════════════

def slide_heaven(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    slide_header(slide, "Heaven Worker Infrastructure",
                 subtitle="Cloudflare Edge · D1 Constitutional Ledger · NOIZYSTREAM v2 · HVS 75/25",
                 icon="⚡", accent=ACCENT_CYAN, slide_num=5)

    # Architecture layers (stacked cards)
    layers = [
        ("EDGE REQUEST",    "noizy.ai/*  →  Cloudflare Worker (Heaven)",
         "CORS · Auth gate (X-Noizy-Key) · Gabriel audit on every write",  ACCENT_CYAN),
        ("DISPATCH LAYER",  "/api/dispatch  →  GOD Mesh (GOD:9696)",
         "CF Access Service Token · MESH_ORIGIN tunnel · 502 fallback with Gabriel error log", ACCENT_GOLD),
        ("STORAGE",         "D1 agent-memory · D1 noizy-prod · D1 aquarium-archive",
         "6 KV namespaces (signups, royalties, guild, sessions, submissions, memcell)", ACCENT_GREEN),
        ("REALTIME",        "NOIZYSTREAM v2 — Durable Objects (SignalingRoom)",
         "WebRTC signaling rooms · /stream/* routes · persistent state per room", ACCENT_PINK),
    ]

    for i, (title, sub, detail, color) in enumerate(layers):
        y = Inches(1.25 + i * 1.38)
        add_rect(slide, Inches(0.45), y, Inches(8.5), Inches(1.25),
                 fill_color=GREY_DARK, line_color=color, line_width_pt=1.2)
        add_rect(slide, Inches(0.45), y, Inches(0.08), Inches(1.25), fill_color=color)
        add_text(slide, title, Inches(0.65), y + Inches(0.1), Inches(7.8), Inches(0.35),
                 font_size=12, bold=True, color=color)
        add_text(slide, sub, Inches(0.65), y + Inches(0.42), Inches(7.8), Inches(0.3),
                 font_size=11, color=WHITE)
        add_text(slide, detail, Inches(0.65), y + Inches(0.72), Inches(7.8), Inches(0.45),
                 font_size=9, color=GREY_MID)

    # Right: security + deployment panel
    add_rect(slide, Inches(9.2), Inches(1.25), Inches(3.7), Inches(5.5),
             fill_color=GREY_DARK, line_color=ACCENT_CYAN, line_width_pt=0.8)
    add_text(slide, "SECURITY", Inches(9.4), Inches(1.35), Inches(3.4), Inches(0.35),
             font_size=11, bold=True, color=ACCENT_CYAN)
    add_divider(slide, Inches(9.4), Inches(1.72), Inches(3.3), color=ACCENT_CYAN)

    sec_items = [
        "X-Noizy-Key header auth on all protected routes",
        "C2PA stamps on every consent + voice event",
        "Gabriel: immutable audit trail in D1",
        "CF Access Service Token for Mesh tunnel",
        "HVS 75/25 locked at protocol level",
        "CORS scoped to noizy.ai only",
        "Durable Object state isolation per room",
    ]
    for j, item in enumerate(sec_items):
        add_text(slide, f"▸  {item}",
                 Inches(9.4), Inches(1.85 + j * 0.6), Inches(3.3), Inches(0.55),
                 font_size=9, color=GREY_MID if j % 2 else WHITE)

    # HVS badge
    add_rect(slide, Inches(9.2), Inches(6.72), Inches(3.7), Inches(0.6),
             fill_color=ACCENT_GOLD)
    add_text(slide, "HVS 75 / 25  ·  PERPETUAL  ·  LOCKED",
             Inches(9.25), Inches(6.8), Inches(3.6), Inches(0.45),
             font_size=12, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — CROSS-PLATFORM FOOTPRINT
# ══════════════════════════════════════════════════════════════════════════════

def slide_platforms(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    slide_header(slide, "Cross-Platform Command Footprint",
                 subtitle="One empire. Every surface. Zero dead zones.",
                 icon="⬢", accent=ACCENT_PINK, slide_num=6)

    platforms = [
        {
            "name":    "M2 ULTRA (GOD)",
            "tooling": "Docker · Swift · FastAPI · uvloop",
            "purpose": "Primary compute node. Mesh GOD:9696. AI inference. "
                       "Voice server. The Gathering RAID. 12TB+ local vault.",
            "color":   ACCENT_GOLD,
            "icon":    "🖥",
        },
        {
            "name":    "iPAD / iPHONE",
            "tooling": "Lucy Swift App · CF Access · WebRTC",
            "purpose": "Mobile command-and-control. Spatial UI. "
                       "Lucy agent interface. Real-time NOIZYSTREAM v2 sessions.",
            "color":   ACCENT_CYAN,
            "icon":    "📱",
        },
        {
            "name":    "VS CODE / INSIDERS",
            "tooling": "TypeScript Extension · MCP Servers · GitHub Copilot",
            "purpose": "Native AI-assisted engineering. MCP: noizylab-workspace, "
                       "noizy-memory, cloudflare, github. Wrangler deploy pipeline.",
            "color":   ACCENT_GREEN,
            "icon":    "⌨",
        },
        {
            "name":    "WEB — noizy.ai PORTAL",
            "tooling": "Heaven Worker · D1 · KV · Durable Objects",
            "purpose": "Public-facing SSO gateway. Admin HUD. Signup capture. "
                       "Gabriel audit viewer. Royalty dashboard. NOIZYSTREAM rooms.",
            "color":   ACCENT_PINK,
            "icon":    "🌐",
        },
        {
            "name":    "CLAUDE DESKTOP / GOD",
            "tooling": "GitKraken MCP · Docker MCP · NOIZY MCP stack",
            "purpose": "Claude on GOD with full NOIZY MCP mesh attached. "
                       "Filesystem access. Agent coordination. Autonomous ops.",
            "color":   ACCENT_CYAN,
            "icon":    "🤖",
        },
        {
            "name":    "GOOGLE ANTIGRAVITY",
            "tooling": "Azure MCP · CloudRun MCP · Built-in tools",
            "purpose": "Cloud infrastructure orchestration. The GOD-mode IDE. "
                       "GCP + Azure management. Strategy synthesis.",
            "color":   ACCENT_GOLD,
            "icon":    "🚀",
        },
    ]

    cols = 3
    card_w = Inches(4.1)
    card_h = Inches(2.35)
    start_x = Inches(0.45)
    start_y = Inches(1.3)

    for i, p in enumerate(platforms):
        col = i % cols
        row = i // cols
        x = start_x + col * Inches(4.35)
        y = start_y + row * Inches(2.55)

        add_rect(slide, x, y, card_w, card_h,
                 fill_color=GREY_DARK, line_color=p["color"], line_width_pt=1.2)
        # top accent
        add_rect(slide, x, y, card_w, Inches(0.07), fill_color=p["color"])
        # icon + name
        add_text(slide, p["icon"], x + Inches(0.15), y + Inches(0.12),
                 Inches(0.45), Inches(0.45), font_size=18)
        add_text(slide, p["name"], x + Inches(0.6), y + Inches(0.14),
                 card_w - Inches(0.7), Inches(0.38),
                 font_size=13, bold=True, color=p["color"])
        # tooling
        add_text(slide, p["tooling"], x + Inches(0.15), y + Inches(0.56),
                 card_w - Inches(0.25), Inches(0.28),
                 font_size=8, bold=True, color=GREY_MID)
        # divider
        add_divider(slide, x + Inches(0.15), y + Inches(0.85),
                    card_w - Inches(0.25), color=p["color"], thickness=0.6)
        # purpose
        add_text(slide, p["purpose"], x + Inches(0.15), y + Inches(0.95),
                 card_w - Inches(0.25), Inches(1.3),
                 font_size=9, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — TOP 25 EXECUTION ROADMAP
# ══════════════════════════════════════════════════════════════════════════════

def slide_roadmap(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    slide_header(slide, "Top 25 Execution Roadmap — April 17 or Never",
                 subtitle="Four phases. Zero excuses. Gabriel auditing every step.",
                 icon="◎", accent=ACCENT_GREEN, slide_num=7)

    phases = [
        {
            "phase": "PHASE 1",
            "title": "IDENTITY",
            "color": ACCENT_GOLD,
            "items": [
                "Fix Cloudflare email: rsp@ → rsplowman@",
                "Verify Google Workspace SSO OAuth 2.0",
                "MFA on all six brand accounts",
                "wrangler login to account 5f36aa97",
                "Secrets audit: 9 keys on Heaven worker",
                "Rotate NOIZY_KEY pre-launch",
            ],
        },
        {
            "phase": "PHASE 2",
            "title": "DEPLOYMENT",
            "color": ACCENT_CYAN,
            "items": [
                "Deploy Heaven Worker (wrangler deploy)",
                "Run D1 migrations on agent-memory",
                "Smoke test POST /api/signup",
                "Register MC96ECO as founding member",
                "Submit founding consent matrix (C2PA)",
                "Fix NOIZYLAB wrangler name collision",
            ],
        },
        {
            "phase": "PHASE 3",
            "title": "INTEGRATION",
            "color": ACCENT_GREEN,
            "items": [
                "Wire n8n webhook → /api/hook/n8n",
                "Connect Linear API → ticket creation",
                "Connect Notion API → wisdom archive",
                "Test /api/dispatch → GOD Mesh tunnel",
                "Enable R2 Voice Vault in dashboard",
                "Mirror MCP config: all 4 clients",
            ],
        },
        {
            "phase": "PHASE 4",
            "title": "SCALING",
            "color": ACCENT_PINK,
            "items": [
                "Launch The Gathering RAID scan",
                "Scaffold github.com/rsp-noizy/the-gathering",
                "Noah GORUNFREE → Cloud Run deploy",
                "GCP the-gathering-493305 comms prefs",
                "Deploy Lucy Swift iOS app to TestFlight",
                "T-24h full system status sweep",
            ],
        },
    ]

    col_w = Inches(3.1)
    for i, ph in enumerate(phases):
        x = Inches(0.45 + i * 3.18)
        color = ph["color"]

        # Phase header
        add_rect(slide, x, Inches(1.3), col_w, Inches(0.55), fill_color=color)
        add_text(slide, ph["phase"], x + Inches(0.1), Inches(1.32),
                 Inches(1.0), Inches(0.28), font_size=8, bold=True, color=BG_DARK)
        add_text(slide, ph["title"], x + Inches(0.1), Inches(1.52),
                 col_w - Inches(0.15), Inches(0.28), font_size=12, bold=True, color=BG_DARK)

        # Items
        for j, item in enumerate(ph["items"]):
            y = Inches(1.92 + j * 0.82)
            add_rect(slide, x, y, col_w, Inches(0.74),
                     fill_color=GREY_DARK, line_color=color if j == 0 else GREY_DARK,
                     line_width_pt=0.5)
            # item number
            add_rect(slide, x, y, Inches(0.28), Inches(0.74), fill_color=color)
            add_text(slide, str(i * 6 + j + 1), x, y, Inches(0.28), Inches(0.74),
                     font_size=9, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
            add_text(slide, item, x + Inches(0.32), y + Inches(0.12),
                     col_w - Inches(0.38), Inches(0.55),
                     font_size=9, color=WHITE)

    # Launch date stamp
    add_rect(slide, Inches(0.45), Inches(7.05), Inches(12.4), Inches(0.35),
             fill_color=ACCENT_GREEN)
    add_text(slide, "🚀  LAUNCH TARGET:  APRIL 17, 2026  —  God Node online. Heaven live. "
             "Gabriel watching. HVS 75/25 locked. The empire is sovereign.",
             Inches(0.65), Inches(7.08), Inches(12.0), Inches(0.28),
             font_size=10, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════

def build_deck():
    print("⚡ NOIZY.AI EMPIRE DECK — Building...")
    prs = new_prs()

    print("  [1/7] Cover slide...")
    slide_cover(prs)

    print("  [2/7] SSO Architecture...")
    slide_sso(prs)

    print("  [3/7] Agentic Mesh / Trinity...")
    slide_mesh(prs)

    print("  [4/7] The Gathering...")
    slide_gathering(prs)

    print("  [5/7] Heaven Infrastructure...")
    slide_heaven(prs)

    print("  [6/7] Cross-Platform Footprint...")
    slide_platforms(prs)

    print("  [7/7] Execution Roadmap...")
    slide_roadmap(prs)

    prs.save(str(OUT_FILE))
    print(f"\n✅  DONE → {OUT_FILE}")
    print(f"    Slides: {len(prs.slides)}")
    print(f"    Size:   {OUT_FILE.stat().st_size / 1024:.1f} KB")
    print("\n    Open in Keynote or PowerPoint for best rendering.")
    print("    Gabriel is watching. HVS 75/25. Perpetual.\n")


if __name__ == "__main__":
    build_deck()
